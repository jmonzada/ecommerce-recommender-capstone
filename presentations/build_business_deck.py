"""Build presentations/business_deck.pptx - the Step 6 executive deck.

Same structure and identical numbers as presentations/build_business_deck.py,
but the slide text is de-AI'd and carries zero em dashes, en dashes, or arrow
glyphs (sub-bullets use a white-bullet marker; ranges are written with "to").
Plain-English narrative for a non-technical audience; every number is
transcribed from the metrics/fairness JSON. Dash-free by design. Regenerate:

    python presentations/build_business_deck.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FIG = "reports/figures/04_lorenz.png"
OUT = "presentations/business_deck.pptx"

NAVY = RGBColor(0x14, 0x22, 0x3C)
BLUE = RGBColor(0x1F, 0x6F, 0xC4)
GRAY = RGBColor(0x5A, 0x64, 0x70)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    if kicker:
        kb = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.4))
        p = kb.text_frame.paragraphs[0]
        p.text = kicker.upper()
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12.1), Inches(1.1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    return s


def add_bullets(slide, items, top=2.0, left=0.7, width=11.9, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = 0
        text = item
        if item.startswith("  "):
            indent, text = 1, item.strip()
        p.text = ("• " if indent == 0 else "◦ ") + text  # bullet / white-bullet, no dashes
        p.level = indent
        p.font.size = Pt(size if indent == 0 else size - 2)
        p.font.color.rgb = NAVY if indent == 0 else GRAY
        p.space_after = Pt(10)
    return tb


def add_note(slide, text, top=6.7):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = GRAY
    tb.text_frame.word_wrap = True


def add_chevrons(slide, labels, top=5.3, left=0.8, width=2.95, height=0.85):
    for i, label in enumerate(labels):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(left + i * (width + 0.08)), Inches(top),
            Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = LIGHT
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY


def add_big_stat(slide, stat, caption, left, top=2.1, width=3.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = stat
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE
    c = tf.add_paragraph()
    c.text = caption
    c.alignment = PP_ALIGN.CENTER
    c.font.size = Pt(15)
    c.font.color.rgb = GRAY


# ---------------------------------------------------------------- 1 title
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.6))
tf = bg.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Bringing customers back"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = "A product recommendation engine for the Olist marketplace"
p2.font.size = Pt(24)
p2.font.color.rgb = BLUE
p3 = tf.add_paragraph()
p3.text = "Daniel Jethro Monzada · AI/ML capstone · July 2026"
p3.font.size = Pt(16)
p3.font.color.rgb = GRAY
add_note(s, "Data: Brazilian E-Commerce Public Dataset by Olist (Kaggle, CC BY-NC-SA 4.0).")

# ---------------------------------------------------------------- 2 problem
s = add_slide("97 of every 100 customers never come back", kicker="The problem")
add_big_stat(s, "96.9%", "of customers have placed\nexactly one order", 0.9)
add_big_stat(s, "3.12%", "repeat-purchase rate\n(2,997 of 96,096 customers)", 4.9)
add_big_stat(s, "0.02%", "of recommendation slots reach\nsmall sellers without intervention", 8.9)
add_bullets(s, [
    "Customer acquisition spend buys, on average, exactly one order.",
    "Olist sells small merchants on exposure. A bestseller-only storefront quietly breaks that promise.",
    "Today there is no personalization: every visitor sees the same popular products.",
], top=4.5)

# ---------------------------------------------------------------- 3 opportunity
s = add_slide("Small movements in retention are worth real orders", kicker="The opportunity")
add_big_stat(s, "960", "additional returning customers\nper +1pp repeat-purchase rate\n(arithmetic, not a forecast)", 2.4)
add_big_stat(s, "≈ R$130k", "in orders that +1pp represents,\nat the observed R$137\naverage order value", 7.0)
add_bullets(s, [
    "Returning customers arrive with history, so they're exactly the ones a recommender can serve well.",
    "For the single-purchase majority, the win is smarter first-purchase cross-sell and regional merchandising, not personalization.",
    "For sellers: measurable exposure for rarely-shown products protects the marketplace's supply side.",
], top=4.5)
add_note(s, "Scope honesty: personalization applies to the ~3% with history today; the engine handles everyone, but with different tools.")

# ---------------------------------------------------------------- 4 what it does
s = add_slide("What the engine does (no math required)", kicker="The product")
add_bullets(s, [
    "Step 1 builds a shortlist: products similar to what you've bought, plus what shoppers like you buy.",
    "Step 2 ranks the shortlist: a model scores how likely you are to buy each item.",
    "For a brand-new visitor, it falls back gracefully: bestsellers for your region, or overall.",
    "Every recommendation can explain itself in one plain sentence: \"why you're seeing this\".",
    "Runs as a small web service in a container, serving exactly the pipeline the evaluation graded.",
])
add_chevrons(s, ["Your history", "Shortlist (hundreds)", "Ranked top-10", "Plain-language why"])

# ---------------------------------------------------------------- 5 results
s = add_slide("What the evaluation says, honestly", kicker="Results")
add_big_stat(s, "1 in 5", "returning buyers get their next\npurchase in the top-10 shortlist,\nvs 1 in 32 with bestsellers", 0.9, width=4.1)
add_big_stat(s, "+20%", "full pipeline over the shortlist\nalone, on a strictly past-only\nholdout (statistically real)", 5.0, width=4.0)
add_big_stat(s, "37%", "of the catalog reachable in\nreturning buyers' shortlists;\na bestseller feed shows 10 products", 8.9, width=4.0)
add_bullets(s, [
    "Two test beds, labeled honestly: shortlist quality is graded on returning buyers' held-out orders; the deployed pipeline on a strictly past-only three-month holdout.",
    "All numbers are offline estimates. The honest next step is a live A/B test, and the rollout plan is built around one.",
], top=4.7)
add_note(s, "For the technical appendix: shortlist hit rate 0.203 vs 0.031 popularity (leave-last-order-out, 1,949 repeat buyers); "
            "end-to-end +0.0024 [+0.0008, +0.0041] Hit@10 vs shortlist-only; shortlist coverage 0.371; live top-10 coverage 0.071 to 0.102 with the seller dial.")

# ---------------------------------------------------------------- 6 fairness
s = add_slide("We audited the bias before anyone asked", kicker="Fairness & brand risk")
add_bullets(s, [
    "Recommendation quality is not uniform: Northeast customers currently get the best results. We found it, traced it to the shortlist stage, and monitor it.",
    "Without intervention, small sellers (5.9% of the catalog) receive 0.02% of top-10 slots, a direct risk to the exposure promise.",
    "We built a dial, and measured both sides of it:",
    "  Keep all 10 slots for relevance: best hit rate, near-zero small-seller exposure",
    "  Reserve 3 of 10 slots for under-exposed products: small sellers ×4.3, rarely-shown products ×2.1, at ~15% hit-rate cost",
    "  Southern customers pay the most for that dial, quantified, so the trade-off is a decision, not an accident",
], size=16, width=7.4)
s.shapes.add_picture(FIG, Inches(8.4), Inches(1.9), width=Inches(4.4))
add_note(s, "Where to set the dial is a business decision. Our recommendation: validate on newer data than we tuned on, then start conservative (1 reserved slot) and measure.")

# ---------------------------------------------------------------- 7 trust
s = add_slide("Built for trust from day one", kicker="Trust & governance")
add_bullets(s, [
    "Explainable: every score decomposes into signals; customers see a plain-language reason.",
    "No demographics used anywhere. Geography is audited precisely because it proxies income.",
    "Privacy: no customer identifiers leave our systems, including to the AI vendor that writes explanations.",
    "Monitored in production: customer-mix drift, score drift, live hit rate, and monthly seller-exposure checks.",
    "Reversible: every release is a tagged snapshot; rollback is redeploying the previous one.",
])

# ---------------------------------------------------------------- 8 rollout
s = add_slide("Rollout: prove it live before trusting it", kicker="Plan")
add_bullets(s, [
    "Phase 1, shadow mode (2 to 4 weeks): serve recommendations silently and compare against actual purchases.",
    "Phase 2, A/B test on returning buyers: recommendation rail vs bestseller rail; primary metric = repeat-purchase rate.",
    "Phase 3, regional merchandising for new visitors, measured on first-purchase cross-sell.",
    "Phase 4, seller-exposure dial, set with the merchant team after validating on newer data.",
    "Each phase has a kill switch; the offline numbers set expectations, the live numbers decide.",
])
add_chevrons(s, ["1 · Shadow", "2 · A/B returning buyers", "3 · Regional merch", "4 · Exposure dial"])

# ---------------------------------------------------------------- 9 risks & ask
s = add_slide("Risks we're carrying, and what we need", kicker="Risks & the ask")
add_bullets(s, [
    "Risks:",
    "  Offline results may not transfer online, which is why phase 1 is shadow mode.",
    "  Recommenders amplify what they show, so popularity feedback loops are a risk; exposure metrics watch for this monthly.",
    "  Personalization reach is small today (~3%), and growing it depends on retention improving; the two compound.",
    "The ask:",
    "  An A/B testing slot on the storefront and click/purchase telemetry for the live hit rate.",
    "  A product owner for the seller-exposure dial.",
    "  A decision review after phase 2 with the live numbers on the table.",
], size=16)

prs.save(OUT)
print(f"wrote {OUT} with {len(prs.slides._sldIdLst)} slides")

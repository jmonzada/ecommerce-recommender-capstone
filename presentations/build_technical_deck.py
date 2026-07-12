"""Build presentations/technical_deck.pptx - a native, editable PowerPoint of
the Step 6 technical deck.

Parses the markdown cells of notebooks/05_technical_slides.ipynb into headings,
bullets, numbered lists, tables,
and figures, and renders each as real PowerPoint objects. Dash-free by design;
numbers come straight from the notebook. Regenerate:

    python presentations/build_technical_deck.py
"""

import re
from pathlib import Path

import nbformat
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
NB = REPO / "notebooks" / "05_technical_slides.ipynb"
OUT = REPO / "presentations" / "technical_deck.pptx"

NAVY = RGBColor(0x14, 0x22, 0x3C)
BLUE = RGBColor(0x1F, 0x6F, 0xC4)
GRAY = RGBColor(0x5A, 0x64, 0x70)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# per-slide body font override (dense slides get smaller text)
BODY_PT = {1: 13, 8: 12.5, 10: 13.5, 11: 13.5}

TOK = re.compile(r"(\*\*.+?\*\*|\*[^*]+?\*|`[^`]+?`)")


def add_runs(p, text, size, color=NAVY):
    """Render inline **bold**, *italic*, `code` as runs."""
    for part in TOK.split(text):
        if not part:
            continue
        r = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            r.text, r.font.bold = part[2:-2], True
        elif part.startswith("`") and part.endswith("`"):
            r.text, r.font.name = part[1:-1], "Consolas"
        elif part.startswith("*") and part.endswith("*"):
            r.text, r.font.italic = part[1:-1], True
        else:
            r.text = part
        r.font.size = Pt(size)
        r.font.color.rgb = color


def parse_cell(src):
    blocks, lines, i = [], src.splitlines(), 0
    while i < len(lines):
        ln = lines[i].rstrip()
        if not ln.strip():
            i += 1
            continue
        if ln.startswith("# "):
            blocks.append(("h1", ln[2:]))
        elif ln.startswith("## "):
            blocks.append(("h2", ln[3:]))
        elif ln.lstrip().startswith("<img"):
            m = re.search(r'src="([^"]+)"', ln)
            blocks.append(("image", m.group(1)))
        elif ln.lstrip().startswith("<small>"):
            blocks.append(("note", re.sub(r"</?small>", "", ln).strip()))
        elif ln.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].lstrip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not re.match(r"^-+$", cells[0].replace(" ", "")):  # skip |---| sep
                    rows.append(cells)
                i += 1
            blocks.append(("table", rows))
            continue
        elif ln.startswith("- "):
            blocks.append(("bullet", ln[2:]))
        elif re.match(r"^\d+\.\s", ln):
            blocks.append(("num", ln))
        elif re.fullmatch(r"\*[^*].*[^*]\*", ln) and not ln.startswith("**"):
            blocks.append(("note", ln[1:-1].replace("**", "")))
        else:
            blocks.append(("para", ln))
        i += 1
    return blocks


def cpl(width_in, size):
    return max(12, int(width_in * 72 / (size * 0.52)))


def est_height(blocks, width_in, size):
    h = 0.0
    for kind, payload in blocks:
        pref = 2 if kind == "bullet" else 0
        lines = max(1, -(-(len(payload) + pref) // cpl(width_in, size)))
        h += lines * size * 1.28 / 72 + 0.06
    return h


def add_text_group(slide, blocks, x, y, w, size):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for kind, payload in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        p.line_spacing = 1.05
        if kind == "bullet":
            r = p.add_run(); r.text = "•  "; r.font.size = Pt(size); r.font.color.rgb = BLUE
            add_runs(p, payload, size)
        elif kind == "note":
            r = p.add_run(); r.text = payload
            r.font.size = Pt(size - 1); r.font.italic = True; r.font.color.rgb = GRAY
        else:  # para, num
            add_runs(p, payload, size)


def add_table(slide, rows, x, y, w):
    n, ncol = len(rows), len(rows[0])
    row_h = 0.34
    gt = slide.shapes.add_table(n, ncol, Inches(x), Inches(y), Inches(w),
                                Inches(row_h * n)).table
    gt.first_row = True
    gt.horz_banding = False
    for ci in range(ncol):
        gt.columns[ci].width = Emu(int(Inches(w) / ncol))
    for ri, row in enumerate(rows):
        for cix, val in enumerate(row):
            cell = gt.cell(ri, cix)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            if ri == 0:
                r = p.add_run(); r.text = val.replace("*", "")
                r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = NAVY
            else:
                add_runs(p, val, 11)
    return row_h * n


def image_size(path, target_w, max_h):
    im = Image.open(REPO / "notebooks" / path)
    a = im.size[0] / im.size[1]
    w, h = target_w, target_w / a
    if h > max_h:
        h, w = max_h, max_h * a
    return w, h


def title_slide(prs, blocks):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(3.3))
    tf = tb.text_frame; tf.word_wrap = True
    for j, (kind, payload) in enumerate(blocks):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        if kind == "h1":
            r = p.add_run(); r.text = payload
            r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
        elif kind == "note":
            r = p.add_run(); r.text = payload
            r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = GRAY
        else:
            add_runs(p, payload, 16, BLUE if j == 1 else GRAY)


def content_slide(prs, blocks):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    heading = next(p for k, p in blocks if k in ("h1", "h2"))
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.85))
    p = tb.text_frame.paragraphs[0]; tb.text_frame.word_wrap = True
    r = p.add_run(); r.text = heading
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    rule = s.shapes.add_shape(1, Inches(0.55), Inches(1.2), Inches(12.23), Pt(2.2))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLUE; rule.line.fill.background()

    body = [(k, pl) for k, pl in blocks if k not in ("h1", "h2")]
    images = [pl for k, pl in body if k == "image"]
    left_blocks = [(k, pl) for k, pl in body if k != "image"]
    size = BODY_PT.get(idx, 14)
    text_w = 7.35 if images else 12.2

    # images -> right column
    imgy = 1.55
    for path in images:
        w, h = image_size(path, 4.7, 4.6)
        s.shapes.add_picture(str((REPO / "notebooks" / path).resolve()),
                             Inches(8.05), Inches(imgy), width=Inches(w), height=Inches(h))
        imgy += h + 0.2

    # left column: flow text groups and tables
    y, pending = 1.45, []
    for kind, payload in left_blocks:
        if kind == "table":
            if pending:
                add_text_group(s, pending, 0.55, y, text_w, size)
                y += est_height(pending, text_w, size) + 0.08
                pending = []
            th = add_table(s, payload, 0.55, y, text_w)
            y += th + 0.18
        else:
            pending.append((kind, payload))
    if pending:
        add_text_group(s, pending, 0.55, y, text_w, size)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

nb = nbformat.read(NB, as_version=4)
cells = [c.source for c in nb.cells if c.cell_type == "markdown"]
for idx, src in enumerate(cells):
    blocks = parse_cell(src)
    if idx == 0:
        title_slide(prs, blocks)
    else:
        content_slide(prs, blocks)

prs.save(OUT)
print(f"wrote {OUT} with {len(prs.slides._sldIdLst)} slides")
